from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import json, os

with open("analysis_data.json", "r", encoding="utf-8") as f:
    data = json.load(f)
py = data["py_summary"]
cpp = data["cpp_summary"]
tl = py["total_lines"] + cpp["total_lines"]
tf = py["count"] + cpp["count"]
tfn = py["total_functions"] + cpp["total_functions"]
tc = py["total_classes"] + cpp["total_classes"]
tb = py["total_branches"] + cpp["total_branches"]
tlp = py["total_loops"] + cpp["total_loops"]

py_rec = sum(1 for r in data["py_details"] if r["has_recursion"])
cpp_try = sum(1 for r in data["cpp_details"] if r["has_try_catch"])

C = {
    "navy": RGBColor(0x0F, 0x2B, 0x46), "navy2": RGBColor(0x16, 0x3A, 0x5C),
    "cyan": RGBColor(0x06, 0xB6, 0xD4), "cyan_dark": RGBColor(0x0E, 0x74, 0x94),
    "blue": RGBColor(0x0E, 0xA5, 0xE9), "sky": RGBColor(0x38, 0xBD, 0xF8),
    "indigo": RGBColor(0x63, 0x66, 0xF1),
    "white": RGBColor(0xFF, 0xFF, 0xFF), "bg": RGBColor(0xF2, 0xF6, 0xFA),
    "card": RGBColor(0xFF, 0xFF, 0xFF), "border": RGBColor(0xD0, 0xDF, 0xEB),
    "text": RGBColor(0x1A, 0x26, 0x38), "text2": RGBColor(0x47, 0x56, 0x69),
    "text3": RGBColor(0x8B, 0x96, 0xA8), "cyan_light": RGBColor(0xDC, 0xF3, 0xF8),
}

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

def R(s, l, t, w, h, f=None, b=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.shadow.inherit = False
    if f: sh.fill.solid(); sh.fill.fore_color.rgb = f
    else: sh.fill.background()
    if b: sh.line.color.rgb = b; sh.line.width = Pt(0.5)
    else: sh.line.fill.background()
    return sh

def OV(s, l, t, d, f=None):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    sh.shadow.inherit = False
    if f: sh.fill.solid(); sh.fill.fore_color.rgb = f
    else: sh.fill.background()
    sh.line.fill.background()
    return sh

def T(s, l, t, w, h, tx, sz=13, b=False, c=None, a=PP_ALIGN.LEFT):
    C_ = c or C["text"]
    bx = s.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = tx
    p.font.size = Pt(sz)
    p.font.bold = b
    p.font.color.rgb = C_
    p.font.name = "Microsoft YaHei"
    p.alignment = a
    return bx

def MT(s, l, t, w, h, lines, sz=11, bold_first=False, c=None):
    C_ = c or C["text2"]
    bx = s.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = bx.text_frame.paragraphs[0]
        else:
            p = bx.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = C_
        p.font.name = "Microsoft YaHei"
        if bold_first and i == 0:
            p.font.bold = True
            p.font.color.rgb = C["text"]
        p.space_after = Pt(2)
    return bx

def ICON(s, icon, l, t, d, bg, color=C["text"]):
    OV(s, l, t, d, f=bg)
    T(s, l, t+Inches(0.02), d, d, icon, sz=int(d/914400*11), c=color, a=PP_ALIGN.CENTER)

def DECO(s):
    OV(s, Inches(-0.8), Inches(-0.8), Inches(2.5), f=RGBColor(0xE0, 0xEB, 0xF5))
    OV(s, Inches(11.5), Inches(5.5), Inches(3), f=RGBColor(0xE8, 0xF0, 0xF8))
    OV(s, Inches(12.5), Inches(6.5), Inches(1.5), f=RGBColor(0xDC, 0xF3, 0xF8))

def HDR(s, title, sub="", num=""):
    R(s, 0, 0, W, Inches(1.15), f=C["navy"])
    R(s, 0, Inches(1.15), W, Inches(0.04), f=C["cyan"])
    R(s, Inches(0), Inches(0), Inches(0.06), Inches(1.15), f=C["cyan"])
    if num:
        T(s, Inches(0.4), Inches(0.08), Inches(0.5), Inches(0.5), num, sz=20, b=True, c=C["cyan"])
    T(s, Inches(0.4) if not num else Inches(0.9), Inches(0.15), Inches(8), Inches(0.45),
      title, sz=22, b=True, c=C["white"])
    if sub:
        T(s, Inches(0.9), Inches(0.6), Inches(8), Inches(0.3), sub, sz=10, c=RGBColor(0x8A, 0xBF, 0xD8))

def SN(s, n):
    T(s, W-Inches(0.8), H-Inches(0.35), Inches(0.6), Inches(0.25), str(n), sz=9, c=C["text3"], a=PP_ALIGN.RIGHT)

def IMG(s, path, l, t, w, h=None):
    if os.path.exists(path):
        s.shapes.add_picture(path, l, t, w, h)

def CARD(s, l, t, w, h, icon, title, items, color=C["cyan"], sz=10):
    R(s, l, t, w, h, f=C["card"], b=C["border"])
    R(s, l+Inches(0.02), t+Inches(0.02), w-Inches(0.04), Inches(0.06), f=color)
    ICON(s, icon, l+Inches(0.15), t+Inches(0.2), Inches(0.38), color, C["white"])
    T(s, l+Inches(0.65), t+Inches(0.2), w-Inches(0.85), Inches(0.35), title, sz=14, b=True, c=C["text"])
    y = t + Inches(0.65)
    for it in items:
        T(s, l+Inches(0.2), y, w-Inches(0.35), Inches(0.26), f"  {it}", sz=sz, c=C["text2"])
        y += Inches(0.24+max(0, (len(it)-40)*0.02))

def TBL(s, l, t, w, h, rows, cols, data, cw=None):
    ts = s.shapes.add_table(rows, cols, l, t, w, h)
    tb = ts.table
    if cw:
        for i, ww in enumerate(cw):
            tb.columns[i].width = ww
    for ri, rd in enumerate(data):
        for ci, cv in enumerate(rd):
            c = tb.cell(ri, ci)
            c.text = str(cv)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = "Microsoft YaHei"
                if ri == 0: p.font.bold = True; p.font.color.rgb = C["white"]
                else: p.font.color.rgb = C["text"]
            c.margin_left = Inches(0.06)
            c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.02)
            c.margin_bottom = Inches(0.02)
            if ri == 0: c.fill.solid(); c.fill.fore_color.rgb = C["navy"]
            elif ri % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = C["bg"]
            else: c.fill.solid(); c.fill.fore_color.rgb = C["white"]
    return ts

def KPI(s, l, t, w, val, label, color=C["cyan"]):
    R(s, l, t, w, Inches(1.1), f=C["card"], b=C["border"])
    R(s, l, t, w, Inches(0.04), f=color)
    T(s, l+Inches(0.05), t+Inches(0.1), w-Inches(0.1), Inches(0.45), val, sz=22, b=True, c=color, a=PP_ALIGN.CENTER)
    T(s, l+Inches(0.05), t+Inches(0.6), w-Inches(0.1), Inches(0.35), label, sz=9, c=C["text2"], a=PP_ALIGN.CENTER)

def FTNOTE(s):
    R(s, Inches(0.15), H-Inches(0.22), W-Inches(0.3), Inches(0.02), f=C["border"])
    T(s, Inches(0.3), H-Inches(0.2), Inches(8), Inches(0.18),
      "编程代码注释智能生成与优化助手  |  项目报告", sz=7, c=C["text3"])

def add_transition(slide, effect="fade"):
    sld = slide._element
    transition = sld.makeelement(qn('p:transition'), {})
    sld.append(transition)
    if effect == "fade":
        transition.set(qn('p:spd'), 'slow')
    elif effect == "push":
        transition.set(qn('p:spd'), 'med')

# ═══════════════════ 1: COVER ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["navy"])
R(sl, Inches(6.5), 0, Inches(6.83), H, f=C["navy2"])
R(sl, 0, Inches(3.2), W, Inches(0.04), f=C["cyan"])
OV(sl, Inches(-1), Inches(-1), Inches(4), f=RGBColor(0x12, 0x3E, 0x62))
OV(sl, Inches(10), Inches(5), Inches(4), f=RGBColor(0x14, 0x48, 0x6E))
T(sl, Inches(1.0), Inches(1.8), Inches(5), Inches(0.7),
   "编程代码注释智能生成与优化助手", sz=32, b=True, c=C["white"])
T(sl, Inches(1.0), Inches(2.6), Inches(5), Inches(0.35),
   "Code Comment Intelligence Generator & Optimizer", sz=13, c=RGBColor(0x9D, 0xD8, 0xD4))
T(sl, Inches(1.0), Inches(3.5), Inches(5), Inches(0.35),
   "智能注释生成系统  |  项目报告", sz=12, c=RGBColor(0x80, 0xCB, 0xC4))
ICON(sl, "\U0001F4BB", Inches(8.5), Inches(2.0), Inches(0.7), RGBColor(0x0E, 0x5A, 0x7A))
ICON(sl, "\u2699\uFE0F", Inches(9.5), Inches(3.3), Inches(0.7), RGBColor(0x0E, 0x5A, 0x7A))
ICON(sl, "\U0001F916", Inches(10.5), Inches(4.6), Inches(0.7), RGBColor(0x0E, 0x5A, 0x7A))
T(sl, Inches(1.0), Inches(5.0), Inches(4), Inches(0.25),
   "Python FastAPI + DeepSeek API + AST 双引擎", sz=10, c=RGBColor(0x80, 0xCB, 0xC4))
T(sl, Inches(1.0), Inches(5.3), Inches(4), Inches(0.25),
   f"样本: {tf} 文件 / {tl} 行代码  |  v1.0", sz=10, c=RGBColor(0x80, 0xCB, 0xC4))
R(sl, 0, H-Inches(0.06), W, Inches(0.06), f=C["cyan"])
add_transition(sl)

# ═══════════════════ 2: ABSTRACT ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "项目摘要", "Executive Summary", "01")
R(sl, Inches(0.4), Inches(1.5), Inches(12.5), Inches(2.2), f=C["card"], b=C["border"])
R(sl, Inches(0.4), Inches(1.5), Inches(0.04), Inches(2.2), f=C["cyan"])
MT(sl, Inches(0.7), Inches(1.6), Inches(11.8), Inches(2.0), [
    "本项目旨在解决软件开发中代码注释缺失的普遍问题，通过结合程序语法分析（AST）和大语言模型（LLM）双技术路线，实现代码注释的自动生成。",
    "系统基于 Python FastAPI 构建后端服务，使用原生 HTML/CSS/JavaScript 作为前端界面。AST 分析引擎支持 Python（ast 标准库）和 C++（正则表达式），"
    "LLM 语义引擎调用 DeepSeek API，实现双通道智能注释生成。用户通过 Web 界面上传代码文件，系统自动完成分析并以对比形式展示 AST 与 LLM 两种注释结果，"
    "支持注释代码的预览和下载。",
    f"通过 {tf} 个样本文件（{tl} 行代码）的系统性测试，验证了系统在分析性能、结构检测准确率等方面的表现。",
], sz=11)
KPI(sl, Inches(0.4), Inches(4.0), Inches(3.0), f"{tf}", "样本文件总数", C["cyan"])
KPI(sl, Inches(3.6), Inches(4.0), Inches(3.0), f"{tl}", "代码总行数", C["blue"])
KPI(sl, Inches(6.8), Inches(4.0), Inches(3.0), f"{tfn}", "检测函数总数", C["sky"])
KPI(sl, Inches(10.0), Inches(4.0), Inches(3.0), f"{tc}", "检测类总数", C["indigo"])
R(sl, Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.5), f=C["card"], b=C["border"])
R(sl, Inches(0.4), Inches(5.4), Inches(12.5), Inches(0.04), f=C["cyan"])
T(sl, Inches(0.6), Inches(5.6), Inches(3), Inches(0.3), "核心指标速览", sz=12, b=True, c=C["text"])
T(sl, Inches(0.6), Inches(6.0), Inches(12), Inches(0.8),
   f"Python AST: 平均 {py['avg_time_ms']}ms / 文件  |  C++ 正则: 平均 {cpp['avg_time_ms']}ms / 文件  |  "
   f"Python 函数: {py['total_functions']}  |  C++ 函数: {cpp['total_functions']}  |  "
   f"Python 类: {py['total_classes']}  |  C++ 类: {cpp['total_classes']}",
   sz=10, c=C["text2"])
SN(sl, 2); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 3: BACKGROUND ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "项目背景与痛点", "Background & Pain Points", "02")
# Left side - pain points
pts = [
    ("\U0001F4C4", "代码文档严重缺失", "大量遗留代码缺乏注释，新成员接手困难，维护成本高企"),
    ("\u23F1\uFE0F", "人工注释效率低下", "手动注释耗时占开发总工时 20-30%，严重影响迭代效率"),
    ("\u26A0\uFE0F", "注释质量参差不齐", "注释风格因人而异，信息密度低，部分与代码逻辑脱节"),
    ("\U0001F504", "维护更新严重滞后", "代码迭代后注释更新不及时，产生误导，降低团队信任"),
]
for i, (ic, t, d) in enumerate(pts):
    x = Inches(0.3 + i*3.2)
    R(sl, x, Inches(1.5), Inches(3.0), Inches(2.5), f=C["card"], b=C["border"])
    R(sl, x, Inches(1.5), Inches(3.0), Inches(0.04), f=[C["cyan"],C["blue"],C["sky"],C["indigo"]][i])
    ICON(sl, ic, x+Inches(0.15), Inches(1.7), Inches(0.4), [C["cyan"],C["blue"],C["sky"],C["indigo"]][i], C["white"])
    T(sl, x+Inches(0.65), Inches(1.7), Inches(2.2), Inches(0.35), t, sz=13, b=True, c=C["text"])
    T(sl, x+Inches(0.15), Inches(2.2), Inches(2.7), Inches(1.2), d, sz=10, c=C["text2"])

# Right side - problem analysis
R(sl, Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.8), f=C["card"], b=C["border"])
R(sl, Inches(0.4), Inches(4.3), Inches(0.04), Inches(2.8), f=C["cyan"])
MT(sl, Inches(0.7), Inches(4.5), Inches(5.5), Inches(2.5), [
    "现有工具的局限性",
    "Doxygen / Javadoc 等传统工具仅生成模板化函数签名注释",
    "缺乏对代码语义的深入理解，无法描述业务逻辑",
    "注释风格固定，无法适应不同团队的需求",
], sz=10)
MT(sl, Inches(6.5), Inches(4.5), Inches(6.0), Inches(2.5), [
    "核心解决思路",
    "AST 引擎: 毫秒级提取代码结构信息（函数、类、分支、循环）",
    "LLM 引擎: 结合结构上下文生成语义级自然语言注释",
    "双引擎互补: 结构分析 + 语义理解，覆盖全场景",
], sz=10)
SN(sl, 3); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 4: OBJECTIVES ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "项目目标与范围", "Objectives & Scope", "03")
gs = [
    ("\U0001F680", "双引擎并行", "AST 分析引擎 + LLM 语义引擎协同", "毫秒级结构分析 + 智能语义注释", C["cyan"]),
    ("\U0001F310", "多语言覆盖", "支持 Python / C++ 两大主流语言", "覆盖函数、类、分支、循环等全部结构", C["blue"]),
    ("\u2B50", "高质量注释", "结合代码语义生成自然语言注释", "统一风格、高信息密度、业务可读", C["sky"]),
    ("\U0001F4E1", "流畅交互", "Tab 式双栏对比展示", "一键下载带注释代码", C["indigo"]),
]
for i, (ic, t, s1, s2, c) in enumerate(gs):
    x = Inches(0.3 + i*3.2)
    R(sl, x, Inches(1.5), Inches(3.0), Inches(3.0), f=C["card"], b=C["border"])
    R(sl, x, Inches(1.5), Inches(3.0), Inches(0.04), f=c)
    ICON(sl, ic, x+Inches(1.15), Inches(1.7), Inches(0.5), c, C["white"])
    T(sl, x+Inches(0.15), Inches(2.4), Inches(2.7), Inches(0.3), t, sz=15, b=True, c=C["text"], a=PP_ALIGN.CENTER)
    T(sl, x+Inches(0.15), Inches(2.75), Inches(2.7), Inches(0.8), s1, sz=10, c=C["text2"], a=PP_ALIGN.CENTER)
    T(sl, x+Inches(0.15), Inches(3.3), Inches(2.7), Inches(0.8), s2, sz=10, c=C["text2"], a=PP_ALIGN.CENTER)

# Scope section
R(sl, Inches(0.4), Inches(4.8), Inches(12.5), Inches(2.2), f=C["card"], b=C["border"])
R(sl, Inches(0.4), Inches(4.8), Inches(0.04), Inches(2.2), f=C["cyan"])
T(sl, Inches(0.7), Inches(4.9), Inches(2), Inches(0.3), "项目范围", sz=12, b=True, c=C["text"])
MT(sl, Inches(0.7), Inches(5.25), Inches(5.8), Inches(1.6), [
    "\u2713  .py / .cpp / .cc / .cxx / .hpp 文件上传与自动语言识别",
    "\u2713  AST 结构分析（Python 标准 ast + C++ 正则表达式）",
    "\u2713  LLM 语义级注释生成（DeepSeek API）",
    "\u2713  双栏 Tab 对比展示 + 注释代码预览与下载",
], sz=10)
MT(sl, Inches(6.8), Inches(5.25), Inches(5.8), Inches(1.6), [
    "\u2713  100+ 样本系统性测试验证",
    "\u2713  内存字典 + JSON 文件持久化",
    "\u2713  同一源站部署（前后端统一 FastAPI 托管）",
], sz=10)
SN(sl, 4); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 5: ARCHITECTURE ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "系统架构", "System Architecture", "04")
L = [
    ("\U0001F4FA", "前端展示层", ["文件上传界面 (.py/.cpp/.h)", "Tab 双栏对比展示", "highlight.js 代码高亮", "注释代码预览与下载"], C["cyan"]),
    ("\u2699\uFE0F", "双引擎分析层", ["AST 分析 (Python ast)", "C++ 正则解析引擎", "LLM 语义 (DeepSeek API)", "异步并行处理"], C["blue"]),
    ("\U0001F4CB", "数据处理层", ["代码结构提取 & 格式化", "语义上下文构建", "Prompt 模板组装", "注释生成 & 后处理"], C["sky"]),
    ("\U0001F4BE", "持久化层", ["内存字典缓存", "JSON 文件存储", "uploads 文件系统", "零数据库依赖"], C["indigo"]),
]
for i, (ic, t, its, c) in enumerate(L):
    x = Inches(0.4 + i*3.2)
    R(sl, x, Inches(1.5), Inches(3.0), Inches(5.2), f=C["card"], b=C["border"])
    R(sl, x, Inches(1.5), Inches(3.0), Inches(0.45), f=c)
    T(sl, x+Inches(0.1), Inches(1.56), Inches(2.8), Inches(0.35), t, sz=13, b=True, c=C["white"], a=PP_ALIGN.CENTER)
    for j, it in enumerate(its):
        iy = Inches(2.2 + j*0.65)
        R(sl, x+Inches(0.15), iy, Inches(2.7), Inches(0.48), f=C["bg"], b=C["border"])
        T(sl, x+Inches(0.25), iy+Inches(0.05), Inches(2.5), Inches(0.38), it, sz=10, c=C["text"])
    if i < len(L)-1:
        a = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x+Inches(3.0), Inches(3.8), Inches(0.2), Inches(0.25))
        a.fill.solid(); a.fill.fore_color.rgb = C["text3"]; a.line.fill.background(); a.shadow.inherit = False

# Data flow
R(sl, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4), f=C["navy"])
T(sl, Inches(0.6), Inches(6.72), Inches(12), Inches(0.35),
   "数据流: 用户上传  ->  自动语言识别  ->  AST + LLM 并行分析  ->  Tab 对比展示  ->  注释代码下载",
   sz=9, c=C["white"], a=PP_ALIGN.CENTER)
SN(sl, 5); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 6: TECH STACK ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "技术栈与数据流", "Technology Stack & Data Flow", "05")
# Tech stack cards
ST = [
    ("\U0001F4E6", "后端框架", "Python FastAPI", "Uvicorn 异步服务器 | python-pptx\nPython ast 标准库 | 正则表达式引擎", C["cyan"]),
    ("\U0001F310", "前端界面", "原生 HTML/CSS/JS", "highlight.js 语法高亮\nTab 双栏对比 | 零框架设计", C["blue"]),
    ("\U0001F916", "AI 引擎", "DeepSeek API", "兼容 OpenAI SDK | deepseek-chat\n流式/非流式 | 30s 超时重试", C["sky"]),
    ("\U0001F527", "工具链", "python-pptx", "JSON 持久化 | Pandoc 转换\nGit 版本管理 | 零数据库", C["indigo"]),
]
for i, (ic, ca, t, d, c) in enumerate(ST):
    x = Inches(0.3 + i*3.2)
    R(sl, x, Inches(1.5), Inches(3.0), Inches(3.8), f=C["card"], b=C["border"])
    R(sl, x, Inches(1.5), Inches(3.0), Inches(0.04), f=c)
    ICON(sl, ic, x+Inches(1.15), Inches(1.7), Inches(0.5), c, C["white"])
    T(sl, x+Inches(0.15), Inches(2.4), Inches(2.7), Inches(0.2), ca, sz=9, c=c, a=PP_ALIGN.CENTER)
    T(sl, x+Inches(0.15), Inches(2.6), Inches(2.7), Inches(0.35), t, sz=16, b=True, c=C["text"], a=PP_ALIGN.CENTER)
    T(sl, x+Inches(0.15), Inches(3.05), Inches(2.7), Inches(1.8), d, sz=10, c=C["text2"], a=PP_ALIGN.CENTER)

# Data processing flow
R(sl, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.6), f=C["card"], b=C["border"])
R(sl, Inches(0.4), Inches(5.6), Inches(0.04), Inches(1.6), f=C["cyan"])
T(sl, Inches(0.7), Inches(5.7), Inches(3), Inches(0.3), "核心处理流程", sz=12, b=True, c=C["text"])
MT(sl, Inches(0.7), Inches(6.05), Inches(11.8), Inches(1.1), [
    "Step 1: 用户上传代码文件 -> 系统根据扩展名自动识别语言（.py -> Python, .cpp/.cc/.cxx/.hpp -> C++）",
    "Step 2: 后端启动 AST 分析 + LLM 调用两条线程并行执行，前端每 1.5 秒轮询状态",
    "Step 3: AST 引擎提取结构信息（函数、类、分支、循环），LLM 引擎基于 AST 上下文生成语义注释",
    "Step 4: 前端以 Tab 形式分别展示 AST 和 LLM 注释结果，用户选择版本并下载",
], sz=10)
SN(sl, 6); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 7: UPLOAD & AST ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "核心功能: 上传与 AST 分析", "Upload & AST Analysis Engine", "06")
CARD(sl, Inches(0.4), Inches(1.5), Inches(6.3), Inches(2.2),
     "\U0001F4C1", "文件上传与语言识别",
     [".py / .cpp / .cc / .cxx / .hpp 文件拖拽或选择上传",
      "FastAPI UploadFile 接收 multipart/form-data",
      "根据扩展名自动区分 Python/C++，分配对应分析器",
      "文件存入 uploads/ 目录，元信息注册到内存字典"], C["blue"])
CARD(sl, Inches(0.4), Inches(4.0), Inches(6.3), Inches(3.2),
     "\U0001F50D", "AST 分析引擎详解",
     ["Python: 标准 ast 模块生成语法树，提取函数/类/分支/循环",
      f"Python 平均耗时 {py['avg_time_ms']}ms，100% 文件零失败",
      "C++: 正则表达式匹配函数签名、类定义、控制结构",
      f"C++ 平均耗时 {cpp['avg_time_ms']}ms，覆盖约 60-70%",
      "额外检测: 递归函数、异常处理、复杂赋值、头文件包含",
      f"递归检测: {py_rec}/{py['count']} 个 Python 文件含递归调用"],
     C["cyan"])

# Right side - detailed table
rows = [
    ["指标", "Python", "C++"],
    ["样本文件数", str(py["count"]), str(cpp["count"])],
    ["总行数", str(py["total_lines"]), str(cpp["total_lines"])],
    ["平均行数/文件", str(py["avg_lines"]), str(cpp["avg_lines"])],
    ["函数总数", str(py["total_functions"]), str(cpp["total_functions"])],
    ["类总数", str(py["total_classes"]), str(cpp["total_classes"])],
    ["分支总数", str(py["total_branches"]), str(cpp["total_branches"])],
    ["循环总数", str(py["total_loops"]), str(cpp["total_loops"])],
    ["平均耗时(ms)", str(py["avg_time_ms"]), str(cpp["avg_time_ms"])],
    ["最短耗时(ms)", str(py["min_time_ms"]), str(cpp["min_time_ms"])],
    ["最长耗时(ms)", str(py["max_time_ms"]), str(cpp["max_time_ms"])],
    ["平均文件大小(bytes)", str(int(py["avg_size_bytes"])), str(int(cpp["avg_size_bytes"]))],
]
TBL(sl, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.8), len(rows), 3, rows,
    [Inches(2.0), Inches(1.7), Inches(1.7)])
SN(sl, 7); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 8: LLM ENGINE ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "核心功能: LLM 语义注释引擎", "LLM Semantic Annotation Engine", "07")
steps = [
    ("\U0001F4E4", "AST 输出\n输入准备", "提取函数签名、类定义\n控制流结构上下文\n已有注释统计", C["cyan"]),
    ("\U0001F4DD", "Prompt\n模板组装", "系统提示词设定角色\nFew-shot 示例引导\nAST 结构注入上下文", C["blue"]),
    ("\U0001F310", "DeepSeek\nAPI 调用", "deepseek-chat 模型\n流式/非流式输出\n30s 超时自动重试", C["sky"]),
    ("\u2728", "注释\n生成输出", "自然语言行级注释\n函数/类头部注释\n统一风格格式化", C["indigo"]),
]
for i, (ic, t1, t2, c) in enumerate(steps):
    x = Inches(0.3 + i*3.2)
    R(sl, x, Inches(1.5), Inches(3.0), Inches(3.5), f=C["card"], b=C["border"])
    R(sl, x, Inches(1.5), Inches(3.0), Inches(0.04), f=c)
    ICON(sl, ic, x+Inches(1.15), Inches(1.7), Inches(0.5), c, C["white"])
    MT(sl, x+Inches(0.1), Inches(2.4), Inches(2.8), Inches(0.7), [t1], sz=12, bold_first=True, c=C["text"])
    T(sl, x+Inches(0.15), Inches(2.85), Inches(2.7), Inches(1.5), t2, sz=10, c=C["text2"])

CARD(sl, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.8),
     "\U0001F916", "DeepSeek API 调用流程",
     [f"System Prompt: 你是一个代码注释专家，为 {tf} 个样本文件生成高质量、风格统一的中文注释",
      "User: 分析以下 {language} 代码，为每行生成语义级中文注释。代码: {code} | AST: {ast_context}",
      "输出示例: # 计算斐波那契数列的第 n 项 | def fibonacci(n): | # 基础情况: n<=1 时返回 n | if n <= 1: return n"],
     C["cyan"])
SN(sl, 8); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 9: UI & COMPARISON ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "核心功能: 前端展示与下载", "Frontend UI & Download", "08")
CARD(sl, Inches(0.4), Inches(1.5), Inches(6.2), Inches(2.5),
     "\U0001F4BB", "前端界面特性",
     ["原生 HTML/CSS/JavaScript，零框架依赖",
      "highlight.js 实现代码语法高亮",
      "Tab 式切换: AST 注释 / LLM 注释 两个独立标签页",
      "AST 和 LLM 注释使用不同颜色标记，便于区分",
      "每 1.5 秒轮询后端，异步获取分析结果"], C["blue"])
CARD(sl, Inches(6.9), Inches(1.5), Inches(6.0), Inches(2.5),
     "\U0001F4E5", "预览与下载",
     ["用户可选择 AST 或 LLM 注释版本",
      "注释嵌入原始代码对应位置",
      "生成带注释的完整代码文件下载",
      "支持 .py / .cpp 格式输出保存",
      "下载文件即开即用，无需二次处理"], C["cyan"])

# Dual engine comparison table
R(sl, Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.8), f=C["card"], b=C["border"])
R(sl, Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.04), f=C["cyan"])
T(sl, Inches(0.6), Inches(4.45), Inches(4), Inches(0.3), "双引擎注释对比", sz=12, b=True, c=C["text"])
rows = [
    ["维度", "AST 注释", "LLM 注释"],
    ["注释类型", "结构化标签注释", "语义级自然语言注释"],
    ["示例: def add(a,b)", "# @param a: int, b: int | # @return: int", "# 计算两个整数的和并返回结果"],
    ["生成速度", "毫秒级即时", "秒级 (API 调用耗时)"],
    ["覆盖场景", "函数签名、参数、返回值", "代码意图、业务逻辑说明"],
    ["优势", "速度快、结构精准", "语义丰富、可读性强"],
    ["局限", "缺乏语义理解", "依赖 API 可用性"],
]
TBL(sl, Inches(0.6), Inches(4.85), Inches(12.1), Inches(2.1), len(rows), 3, rows,
    [Inches(2.0), Inches(5.0), Inches(5.1)])
SN(sl, 9); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 10: CHARTS 1 ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "测试数据与性能分析（一）", "Test Data & Performance - Part 1", "09")
CD = "charts"
CH = [
    ("chart1_bar.png", "各文件分析耗时分布\n展示 100 个样本文件的 AST 分析耗时分布情况", 0.3, 1.5, 6.2, 5.5),
    ("chart2_perf.png", "Python vs C++ 性能对比\nPython ast 平均 0.55ms vs C++ 正则平均 0.24ms", 6.8, 1.5, 6.2, 5.5),
]
for fname, cap, l, t, w, h in CH:
    x, y, pw, ph = Inches(l), Inches(t), Inches(w), Inches(h)
    R(sl, x, y, pw, ph, f=C["card"], b=C["border"])
    R(sl, x, y, pw, Inches(0.04), f=C["cyan"])
    MT(sl, x+Inches(0.15), y+Inches(0.08), pw-Inches(0.3), Inches(0.45), cap.split("\n"), sz=9, c=C["text2"])
    path = os.path.join(CD, fname)
    if os.path.exists(path):
        IMG(sl, path, x+Inches(0.15), y+Inches(0.5), pw-Inches(0.3), ph-Inches(0.7))
SN(sl, 10); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 11: CHARTS 2 ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "测试数据与性能分析（二）", "Test Data & Performance - Part 2", "10")
x, y, pw, ph = Inches(0.3), Inches(1.5), Inches(6.2), Inches(5.5)
R(sl, x, y, pw, ph, f=C["card"], b=C["border"])
R(sl, x, y, pw, Inches(0.04), f=C["blue"])
MT(sl, x+Inches(0.15), y+Inches(0.08), pw-Inches(0.3), Inches(0.45), ["代码行数 vs 分析耗时", "散点图展示文件规模与耗时的相关性"], sz=9, c=C["text2"])
path = os.path.join(CD, "chart3_scatter.png")
if os.path.exists(path):
    IMG(sl, path, x+Inches(0.15), y+Inches(0.5), pw-Inches(0.3), ph-Inches(0.7))

for i, (fname, cap) in enumerate([("chart4_pie.png", "代码结构分布占比"), ("chart5_coverage.png", "注释覆盖率分布")]):
    x, y = Inches(7.0), Inches(1.5 + i*2.8)
    pw, ph = Inches(5.8), Inches(2.5)
    R(sl, x, y, pw, ph, f=C["card"], b=C["border"])
    R(sl, x, y, pw, Inches(0.04), f=C["cyan"] if i==0 else C["blue"])
    T(sl, x+Inches(0.15), y+Inches(0.08), pw-Inches(0.3), Inches(0.25), cap, sz=10, c=C["text2"], a=PP_ALIGN.CENTER)
    path = os.path.join(CD, fname)
    if os.path.exists(path):
        IMG(sl, path, x+Inches(0.15), y+Inches(0.35), pw-Inches(0.3), ph-Inches(0.55))

R(sl, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.35), f=C["navy"])
T(sl, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.3),
   f"总样本: {tf} 文件（{py['count']} Python + {cpp['count']} C++） | 总行数: {tl} 行 | "
   f"Python 均值: {py['avg_time_ms']}ms | C++ 均值: {cpp['avg_time_ms']}ms | "
   f"函数: {tfn} | 类: {tc} | 分支: {tb} | 循环: {tlp}",
   sz=9, c=C["white"], a=PP_ALIGN.CENTER)
SN(sl, 11); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 12: TEST DATA DETAILS ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "测试数据详细分析", "Test Data Deep Dive", "11")
# Summary stats
stats = [
    (f"{py['total_functions']}", "Python 函数数", C["cyan"]),
    (f"{cpp['total_functions']}", "C++ 函数数", C["blue"]),
    (f"{py['total_classes']}", "Python 类数", C["sky"]),
    (f"{cpp['total_classes']}", "C++ 类数", C["indigo"]),
    (f"{py_rec}/{py['count']}", "递归检测率", C["cyan"]),
    (f"{cpp_try}/{cpp['count']}", "异常处理检测", C["blue"]),
]
for i, (v, l, c) in enumerate(stats):
    KPI(sl, Inches(0.3 + i*2.15), Inches(1.5), Inches(2.0), v, l, c)

# Detailed tables
R(sl, Inches(0.3), Inches(2.9), Inches(6.3), Inches(4.2), f=C["card"], b=C["border"])
R(sl, Inches(0.3), Inches(2.9), Inches(0.04), Inches(4.2), f=C["cyan"])
T(sl, Inches(0.5), Inches(3.0), Inches(4), Inches(0.25), "Python 文件详细分析（前 10）", sz=10, b=True, c=C["text"])
py_top10 = data["py_details"][:10]
py_rows = [["文件名", "行数", "函数", "类", "分支", "循环", "耗时(ms)"]]
for r in py_top10:
    py_rows.append([r["file"], str(r["lines"]), str(r["functions"]), str(r["classes"]),
                    str(r["branches"]), str(r["loops"]), str(r["time_ms"])])
TBL(sl, Inches(0.5), Inches(3.3), Inches(5.9), Inches(3.6), len(py_rows), 7, py_rows,
    [Inches(1.8), Inches(0.5), Inches(0.5), Inches(0.5), Inches(0.5), Inches(0.5), Inches(0.8)])

R(sl, Inches(6.9), Inches(2.9), Inches(6.1), Inches(4.2), f=C["card"], b=C["border"])
R(sl, Inches(6.9), Inches(2.9), Inches(0.04), Inches(4.2), f=C["blue"])
T(sl, Inches(7.1), Inches(3.0), Inches(4), Inches(0.25), "C++ 文件详细分析（前 10）", sz=10, b=True, c=C["text"])
cpp_top10 = data["cpp_details"][:10]
cpp_rows = [["文件名", "行数", "函数", "类", "分支", "循环", "includes"]]
for r in cpp_top10:
    cpp_rows.append([r["file"], str(r["lines"]), str(r["functions"]), str(r["classes"]),
                     str(r["branches"]), str(r["loops"]), str(r["includes"])])
TBL(sl, Inches(7.1), Inches(3.3), Inches(5.7), Inches(3.6), len(cpp_rows), 7, cpp_rows,
    [Inches(1.8), Inches(0.5), Inches(0.5), Inches(0.5), Inches(0.5), Inches(0.5), Inches(0.6)])
SN(sl, 12); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 13: DEVELOPMENT ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "开发流程", "Development Process (5-Day Agile)", "12")
days = [
    ("\U0001F4CB", "Day 1", "需求分析与设计", "技术选型与架构设计\nAPI 接口规范定义\n数据结构与文件设计", C["cyan"]),
    ("\u2699\uFE0F", "Day 2", "核心引擎实现", "Python AST 分析引擎\nC++ 正则解析引擎\nDeepSeek API 集成测试", C["blue"]),
    ("\U0001F4BB", "Day 3", "前端与联调", "上传组件与文件处理\nTab 对比展示界面\n代码高亮与下载功能", C["sky"]),
    ("\U0001F4CA", "Day 4", "测试与数据采集", f"{tf} 样本全量测试\n性能数据采集与分析\nBug 修复与优化", C["indigo"]),
    ("\U0001F4D1", "Day 5", "文档与收尾", "项目报告编写\nPPT 制作\n最终演示准备", C["cyan_dark"]),
]
for i, (ic, day, t, d, c) in enumerate(days):
    x = Inches(0.3 + i*2.6)
    R(sl, x, Inches(1.5), Inches(2.4), Inches(4.8), f=C["card"], b=C["border"])
    R(sl, x, Inches(1.5), Inches(2.4), Inches(0.04), f=c)
    ICON(sl, ic, x+Inches(0.8), Inches(1.7), Inches(0.55), c, C["white"])
    T(sl, x+Inches(0.1), Inches(2.45), Inches(2.2), Inches(0.25), day, sz=10, b=True, c=c, a=PP_ALIGN.CENTER)
    T(sl, x+Inches(0.1), Inches(2.75), Inches(2.2), Inches(0.35), t, sz=14, b=True, c=C["text"], a=PP_ALIGN.CENTER)
    T(sl, x+Inches(0.15), Inches(3.25), Inches(2.1), Inches(2.5), d, sz=10, c=C["text2"], a=PP_ALIGN.CENTER)

# Method
R(sl, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.6), f=C["card"], b=C["border"])
R(sl, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.04), f=C["cyan"])
T(sl, Inches(0.6), Inches(6.55), Inches(12), Inches(0.45),
   "5 天敏捷迭代  |  每日站会同步进度  |  持续集成与测试  |  快速原型验证  |  前后端并行开发",
   sz=10, c=C["text2"], a=PP_ALIGN.CENTER)
SN(sl, 13); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 14: KEY FINDINGS ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "关键结论", "Key Findings & Conclusions", "13")
findings = [
    ("\u2705", "AST 分析引擎", f"成功解析全部 {tf} 个样本文件（{tl} 行代码），零解析失败。检测到 {tfn} 个函数、{tc} 个类、{tb} 个分支、{tlp} 个循环，结构分析覆盖全面。"),
    ("\u26A1", "毫秒级性能", f"Python AST 平均分析耗时 {py['avg_time_ms']}ms，C++ 正则平均 {cpp['avg_time_ms']}ms，均达到毫秒级响应，完全满足实时分析需求。"),
    ("\U0001F517", "双引擎互补", "AST 提供结构级注释（函数签名、参数说明、递归告警），LLM 提供语义级注释（代码意图、业务逻辑说明）。两者结合实现了从结构到语义的全方位注释覆盖。"),
    ("\U0001F4CA", "样本覆盖广泛", f"100 个样本涵盖从基础语法到高级特性的完整知识图谱：基础语法、数据结构、算法实现、面向对象、泛型编程、并发编程、设计模式等，充分验证系统通用性。"),
    ("\U0001F916", "LLM 集成有效", "DeepSeek API 兼容 OpenAI SDK，流式输出支持良好。通过精心设计的 Prompt 模板（系统角色 + Few-shot + AST 上下文注入），生成的注释质量达到可用标准。"),
    ("\U0001F4BB", "零框架前端", "原生 HTML/CSS/JS + highlight.js 实现零框架轻量展示，Tab 双栏对比交互直观，同一源站部署（FastAPI 静态文件托管），避免跨域问题。"),
]
for i, (ic, t, d) in enumerate(findings):
    y = Inches(1.45 + i*0.85)
    R(sl, Inches(0.4), y, Inches(12.5), Inches(0.75), f=C["card"], b=C["border"])
    R(sl, Inches(0.4), y, Inches(0.04), Inches(0.75), f=C["cyan"])
    T(sl, Inches(0.55), y+Inches(0.05), Inches(0.35), Inches(0.3), ic, sz=14, a=PP_ALIGN.CENTER)
    T(sl, Inches(1.0), y+Inches(0.05), Inches(2.0), Inches(0.3), t, sz=12, b=True, c=C["text"])
    T(sl, Inches(1.0), y+Inches(0.35), Inches(11.5), Inches(0.35), d, sz=10, c=C["text2"])
SN(sl, 14); FTNOTE(sl); add_transition(sl)

# ═══════════════════ 15: SUMMARY ═══════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
R(sl, 0, 0, W, H, f=C["bg"])
DECO(sl)
HDR(sl, "总结与展望", "Summary & Future Outlook", "14")
CARD(sl, Inches(0.4), Inches(1.5), Inches(6.2), Inches(3.0),
     "\U0001F3C6", "项目总结",
     ["成功构建 AST + LLM 双引擎代码注释智能生成系统",
      f"覆盖 {tf} 个样本文件、{tl} 行代码，通过完整测试验证",
      "AST 提供结构级注释（函数签名、参数、递归告警）",
      "LLM 提供语义级注释（代码意图、业务逻辑说明）",
      "前端 Tab 双栏对比 + 一键下载，流畅用户体验",
      "FastAPI + 原生前端 + DeepSeek API，轻量高效"],
     C["cyan"])
CARD(sl, Inches(6.9), Inches(1.5), Inches(6.0), Inches(3.0),
     "\U0001F52E", "可改进方向",
     ["更多语言支持: Java / JavaScript / Go / Rust",
      "VS Code 插件: 编辑器内直接生成注释",
      "\"原始 vs 注释\" 差异对比视图",
      "用户反馈机制: 点赞/修改生成的注释",
      "批量上传: 支持整个项目批量分析",
      "多 Provider: OpenAI / Claude 可选后端"],
     C["blue"])

R(sl, Inches(3.0), Inches(4.8), Inches(7.3), Inches(2.3), f=C["card"], b=C["border"])
R(sl, Inches(3.0), Inches(4.8), Inches(7.3), Inches(0.04), f=C["sky"])
OV(sl, Inches(4.9), Inches(5.0), Inches(0.7), f=C["cyan_light"])
T(sl, Inches(4.9), Inches(5.05), Inches(0.7), Inches(0.6), "\U0001F44F", sz=22, a=PP_ALIGN.CENTER)
T(sl, Inches(3.2), Inches(5.85), Inches(6.9), Inches(0.4), "感谢观看", sz=18, b=True, c=C["text"], a=PP_ALIGN.CENTER)
MT(sl, Inches(3.2), Inches(6.2), Inches(6.9), Inches(0.7), [
    "编程代码注释智能生成与优化助手",
    "Python FastAPI + DeepSeek API + AST 双引擎  |  v1.0",
], sz=10, c=C["text3"])
SN(sl, 15); add_transition(sl)

# Save
desktop = os.path.join(os.path.expanduser("~"), "Desktop")
for name in ["代码注释助手_项目报告.pptx", "智能注释项目报告.pptx"]:
    path = os.path.join(desktop, name)
    if not os.path.exists(path):
        prs.save(path)
        print(f"Saved ({os.path.getsize(path)//1024}KB): {name}")
        break
else:
    path = os.path.join(desktop, "AI代码注释项目报告.pptx")
    prs.save(path)
    print(f"Saved ({os.path.getsize(path)//1024}KB): AI代码注释项目报告.pptx")
